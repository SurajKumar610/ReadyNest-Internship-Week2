import os
from sqlalchemy.orm import Session
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill

from ..models import Project, Dataset, AnalyticsResult, CustomerSegment, Recommendation, Prediction

def generate_pdf_report(project_id: int, dataset_id: int, file_path: str, db: Session):
    """
    Generates a professional corporate PDF report with KPIs, segments, and recommendations.
    """
    project = db.query(Project).filter(Project.id == project_id).first()
    
    # Get KPIs
    kpi_result = db.query(AnalyticsResult).filter(
        AnalyticsResult.dataset_id == dataset_id,
        AnalyticsResult.type == "kpis"
    ).first()
    kpis = kpi_result.results if kpi_result else {}
    
    # Get Recommendations
    recs = db.query(Recommendation).filter(Recommendation.dataset_id == dataset_id).all()
    
    doc = SimpleDocTemplate(file_path, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=26,
        textColor=colors.HexColor('#0F172A'),
        spaceAfter=15
    )
    section_style = ParagraphStyle(
        'SectionHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=colors.HexColor('#1E293B'),
        spaceBefore=15,
        spaceAfter=10
    )
    body_style = ParagraphStyle(
        'ReportBody',
        parent=styles['Normal'],
        fontSize=11,
        textColor=colors.HexColor('#334155'),
        spaceAfter=8
    )
    
    story = []
    
    # Title
    story.append(Paragraph("SIGHTFILL EXECUTIVE SUMMARY REPORT", title_style))
    story.append(Paragraph(f"Project: {project.name}", styles['Normal']))
    story.append(Paragraph(f"Generated: {os.path.basename(file_path)}", styles['Normal']))
    story.append(Spacer(1, 20))
    
    # Section: KPIs
    story.append(Paragraph("Key Performance Indicators (KPIs)", section_style))
    story.append(Spacer(1, 10))
    
    kpi_data = [
        ["KPI Metric", "Value"],
        ["Total Revenue", f"${kpis.get('total_revenue', 0.0):,.2f}"],
        ["Net Profit", f"${kpis.get('total_profit', 0.0):,.2f}"],
        ["Total Orders Placed", f"{kpis.get('total_orders', 0):,}"],
        ["Total Customers Active", f"{kpis.get('total_customers', 0):,}"],
        ["Average Order Value (AOV)", f"${kpis.get('average_order_value', 0.0):,.2f}"],
        ["Retention Rate", f"{kpis.get('retention_rate', 0.0)}%"],
        ["Estimated Churn Rate", f"{kpis.get('churn_rate', 0.0)}%"]
    ]
    
    t = Table(kpi_data, colWidths=[200, 150])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1E293B')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 8),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F8FAFC')),
        ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#E2E8F0')),
        ('PADDING', (0,0), (-1,-1), 6)
    ]))
    
    story.append(t)
    story.append(Spacer(1, 25))
    
    # Section: Recommendations
    story.append(Paragraph("Business Insights & Strategic Recommendations", section_style))
    story.append(Spacer(1, 10))
    
    if recs:
        for idx, rec in enumerate(recs):
            rec_text = f"<b>{idx+1}. [{rec.category.upper()}]</b> - Impact: <b>{rec.impact_score}</b> (Conf: {int((rec.confidence_score or 0.8)*100)}%)<br/>{rec.content}"
            story.append(Paragraph(rec_text, body_style))
            story.append(Spacer(1, 8))
    else:
        story.append(Paragraph("No recommendations generated yet. Run pipeline analysis to compile insights.", body_style))
        
    doc.build(story)
    print(f"PDF Report generated at: {file_path}")

def generate_pptx_report(project_id: int, dataset_id: int, file_path: str, db: Session):
    """
    Generates a professional PowerPoint presentation deck with project findings.
    """
    project = db.query(Project).filter(Project.id == project_id).first()
    kpi_result = db.query(AnalyticsResult).filter(
        AnalyticsResult.dataset_id == dataset_id,
        AnalyticsResult.type == "kpis"
    ).first()
    kpis = kpi_result.results if kpi_result else {}
    
    recs = db.query(Recommendation).filter(Recommendation.dataset_id == dataset_id).limit(4).all()
    
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # title layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SIGHTFILL EXECUTIVE DECK"
    subtitle.text = f"Project: {project.name}\nAutomated Business Intelligence & Forecasting"
    
    # Slide 2: KPI Dashboard
    slide_layout = prs.slide_layouts[5] # blank with title
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Performance Indicators (KPIs)"
    
    # Add KPI details
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(4.5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    tf.text = f"Financial Metrics Overview:"
    p = tf.add_paragraph()
    p.text = f"• Total Revenue: ${kpis.get('total_revenue', 0.0):,.2f}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"• Net Profit: ${kpis.get('total_profit', 0.0):,.2f}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"• Average Order Value: ${kpis.get('average_order_value', 0.0):,.2f}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Customer Growth Metrics:"
    p.space_before = Pt(14)
    p = tf.add_paragraph()
    p.text = f"• Active Customer Base: {kpis.get('total_customers', 0):,}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"• Customer Retention Rate: {kpis.get('retention_rate', 0.0)}%"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"• Customer Churn Probability: {kpis.get('churn_rate', 0.0)}%"
    p.level = 1
    
    # Slide 3: Recommendations Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Strategic Recommendations"
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    
    if recs:
        for rec in recs:
            p = tf.add_paragraph()
            p.text = f"• [{rec.category.upper()}] - {rec.content} (Impact: {rec.impact_score})"
            p.space_after = Pt(8)
    else:
        p = tf.add_paragraph()
        p.text = "No recommendations generated yet. Run pipeline analysis to populate."
        
    prs.save(file_path)
    print(f"PPTX Report generated at: {file_path}")

def generate_xlsx_report(project_id: int, dataset_id: int, file_path: str, db: Session):
    """
    Generates a structured multi-tab Excel sheet.
    """
    wb = openpyxl.Workbook()
    
    # Sheet 1: KPIs
    ws_kpi = wb.active
    ws_kpi.title = "KPI Summary"
    ws_kpi.views.sheetView[0].showGridLines = True
    
    kpi_result = db.query(AnalyticsResult).filter(
        AnalyticsResult.dataset_id == dataset_id,
        AnalyticsResult.type == "kpis"
    ).first()
    kpis = kpi_result.results if kpi_result else {}
    
    ws_kpi['A1'] = "Sightfill KPI Summary Report"
    ws_kpi['A1'].font = Font(name="Calibri", size=16, bold=True)
    ws_kpi.row_dimensions[1].height = 30
    
    headers = ["KPI Metric", "Value"]
    ws_kpi.append([]) # blank
    ws_kpi.append(headers)
    
    # Styles
    header_fill = PatternFill(start_color="1F2937", end_color="1F2937", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    
    ws_kpi['A3'].fill = header_fill
    ws_kpi['A3'].font = header_font
    ws_kpi['B3'].fill = header_fill
    ws_kpi['B3'].font = header_font
    
    kpi_rows = [
        ["Total Revenue", kpis.get("total_revenue", 0.0)],
        ["Net Profit", kpis.get("total_profit", 0.0)],
        ["Total Orders", kpis.get("total_orders", 0)],
        ["Total Customers", kpis.get("total_customers", 0)],
        ["Average Order Value", kpis.get("average_order_value", 0.0)],
        ["Customer Lifetime Value", kpis.get("customer_lifetime_value", 0.0)],
        ["Retention Rate (%)", kpis.get("retention_rate", 0.0)],
        ["Churn Rate (%)", kpis.get("churn_rate", 0.0)]
    ]
    
    for r in kpi_rows:
        ws_kpi.append(r)
        
    # Formatting numbers
    for row in range(4, 12):
        cell_a = ws_kpi.cell(row=row, column=1)
        cell_b = ws_kpi.cell(row=row, column=2)
        cell_a.alignment = Alignment(horizontal="left")
        cell_b.alignment = Alignment(horizontal="right")
        
        # currency format
        if row in [4, 5, 8, 9]:
            cell_b.number_format = '$#,##0.00'
        # integer format
        elif row in [6, 7]:
            cell_b.number_format = '#,##0'
        # percentage format
        elif row in [10, 11]:
            cell_b.number_format = '0.00'
            
    # Sheet 2: Segments
    ws_seg = wb.create_sheet(title="Customer Segments")
    ws_seg.views.sheetView[0].showGridLines = True
    ws_seg.append(["Customer ID", "Customer Name", "Recency", "Frequency", "Monetary Value", "Segment Label"])
    
    # Styling headers
    for col in range(1, 7):
        cell = ws_seg.cell(row=1, column=col)
        cell.fill = header_fill
        cell.font = header_font
        
    segments = db.query(CustomerSegment).filter(CustomerSegment.dataset_id == dataset_id).all()
    for seg in segments:
        ws_seg.append([
            seg.customer_id,
            seg.customer_name or "",
            seg.rfm_recency or 0,
            seg.rfm_frequency or 0,
            seg.rfm_monetary or 0.0,
            seg.segment_name
        ])
        
    # Auto-adjust column widths
    for ws in [ws_kpi, ws_seg]:
        for col in ws.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            col_letter = openpyxl.utils.get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
            
    wb.save(file_path)
    print(f"Excel Report generated at: {file_path}")
